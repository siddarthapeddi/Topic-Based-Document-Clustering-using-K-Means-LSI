"""
File handler module for extracting text from various document types
Supports: PDF, DOCX, TXT, DOC
"""

import os
import tempfile
from werkzeug.utils import secure_filename

# Try to import optional dependencies
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import python_pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'doc', 'pptx'}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB

def allowed_file(filename):
    """Check if file extension is allowed"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    if not HAS_PDF:
        raise ImportError("PyPDF2 is not installed. Install with: pip install PyPDF2")
    
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n"
    except Exception as e:
        raise ValueError(f"Error reading PDF: {str(e)}")
    
    return text.strip()

def extract_text_from_docx(file_path):
    """Extract text from DOCX file"""
    if not HAS_DOCX:
        raise ImportError("python-docx is not installed. Install with: pip install python-docx")
    
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + "\n"
    except Exception as e:
        raise ValueError(f"Error reading DOCX: {str(e)}")
    
    return text.strip()

def extract_text_from_pptx(file_path):
    """Extract text from PPTX file"""
    if not HAS_PPTX:
        raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")
    
    text = ""
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            text += f"--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
    except Exception as e:
        raise ValueError(f"Error reading PPTX: {str(e)}")
    
    return text.strip()

def extract_text_from_txt(file_path):
    """Extract text from TXT file"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
    except UnicodeDecodeError:
        # Try different encoding
        with open(file_path, 'r', encoding='latin-1') as file:
            text = file.read()
    except Exception as e:
        raise ValueError(f"Error reading TXT: {str(e)}")
    
    return text.strip()

def extract_text_from_file(file_path, file_type):
    """
    Extract text from various file types
    
    Args:
        file_path (str): Path to the file
        file_type (str): File extension (pdf, docx, txt, pptx, doc)
    
    Returns:
        str: Extracted text
    """
    file_type = file_type.lower()
    
    if file_type == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_type == 'docx':
        return extract_text_from_docx(file_path)
    elif file_type == 'pptx':
        return extract_text_from_pptx(file_path)
    elif file_type in ['txt', 'doc']:
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

def process_uploaded_files(files):
    """
    Process multiple uploaded files and extract text
    
    Args:
        files: List of file objects from Flask
    
    Returns:
        list: List of extracted documents
    """
    documents = []
    errors = []
    
    for file in files:
        if file.filename == '':
            errors.append("Empty filename")
            continue
        
        if not allowed_file(file.filename):
            errors.append(f"File type not allowed: {file.filename}")
            continue
        
        if len(file.read()) > MAX_FILE_SIZE:
            errors.append(f"File too large: {file.filename} (max 10MB)")
            file.seek(0)
            continue
        
        file.seek(0)  # Reset file pointer
        
        try:
            # Save to temporary file
            with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
                file.save(tmp_file.name)
                
                # Extract file extension
                file_ext = file.filename.rsplit('.', 1)[1].lower()
                
                # Extract text
                text = extract_text_from_file(tmp_file.name, file_ext)
                
                if text:
                    documents.append(text)
                else:
                    errors.append(f"No text extracted from: {file.filename}")
                
                # Clean up
                os.unlink(tmp_file.name)
        
        except ImportError as e:
            errors.append(f"Missing dependency for {file.filename}: {str(e)}")
        except Exception as e:
            errors.append(f"Error processing {file.filename}: {str(e)}")
    
    return documents, errors

def get_supported_formats():
    """Get list of supported file formats with availability status"""
    formats = {
        'txt': True,
        'pdf': HAS_PDF,
        'docx': HAS_DOCX,
        'pptx': HAS_PPTX,
    }
    return formats
